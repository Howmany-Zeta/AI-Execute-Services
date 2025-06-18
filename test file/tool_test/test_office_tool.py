import pytest
import asyncio
import os
import tempfile
import shutil
from unittest.mock import patch, MagicMock
from pathlib import Path

from app.tools.office_tool import (
    OfficeTool,
    OfficeToolError,
    InputValidationError,
    FileOperationError,
    SecurityError,
    ContentValidationError,
    OfficeSettings,
    validate_document
)

# Fixtures
@pytest.fixture
def office_tool():
    """Create an OfficeTool instance for testing."""
    tool = OfficeTool({
        "max_file_size_mb": 10,
        "cache_ttl_seconds": 60,
        "cache_max_items": 10,
        "threadpool_workers": 2,
        "default_font": "Arial",
        "default_font_size": 12
    })
    return tool

@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    temp_dir = tempfile.mkdtemp()
    yield temp_dir
    # Clean up after test
    shutil.rmtree(temp_dir)

@pytest.fixture
def sample_docx_file(temp_dir):
    """Create a sample DOCX file for testing."""
    from docx import Document
    
    file_path = os.path.join(temp_dir, "sample.docx")
    doc = Document()
    doc.add_paragraph("This is a test document.")
    doc.add_paragraph("It has multiple paragraphs.")
    
    # Add a table
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Cell 1"
    table.cell(0, 1).text = "Cell 2"
    table.cell(1, 0).text = "Cell 3"
    table.cell(1, 1).text = "Cell 4"
    
    doc.save(file_path)
    return file_path

@pytest.fixture
def sample_pptx_file(temp_dir):
    """Create a sample PPTX file for testing."""
    from pptx import Presentation
    
    file_path = os.path.join(temp_dir, "sample.pptx")
    prs = Presentation()
    
    # Add a slide with text
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "This is a test slide."
    
    # Add another slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Second Slide"
    slide.placeholders[1].text = "This is another test slide."
    
    prs.save(file_path)
    return file_path

@pytest.fixture
def sample_xlsx_file(temp_dir):
    """Create a sample XLSX file for testing."""
    import pandas as pd
    
    file_path = os.path.join(temp_dir, "sample.xlsx")
    df = pd.DataFrame({
        'Name': ['Alice', 'Bob', 'Charlie'],
        'Age': [25, 30, 35],
        'City': ['New York', 'London', 'Paris']
    })
    df.to_excel(file_path, index=False)
    return file_path

@pytest.fixture
def sample_pdf_file(temp_dir):
    """Create a sample PDF file for testing."""
    # This is a simple text file with .pdf extension for testing
    # In a real scenario, you would use a PDF library to create a proper PDF
    file_path = os.path.join(temp_dir, "sample.pdf")
    with open(file_path, 'w') as f:
        f.write("This is a test PDF file content.")
    return file_path

# Basic functionality tests
@pytest.mark.asyncio
async def test_read_docx(office_tool, sample_docx_file):
    """Test reading content from a DOCX file."""
    result = await office_tool.run('read_docx', file_path=sample_docx_file)
    
    assert isinstance(result, dict)
    assert 'paragraphs' in result
    assert len(result['paragraphs']) == 2
    assert result['paragraphs'][0] == "This is a test document."
    assert result['paragraphs'][1] == "It has multiple paragraphs."
    assert result['tables'] is None  # Tables not included by default

@pytest.mark.asyncio
async def test_read_docx_with_tables(office_tool, sample_docx_file):
    """Test reading content from a DOCX file including tables."""
    result = await office_tool.run('read_docx', file_path=sample_docx_file, include_tables=True)
    
    assert isinstance(result, dict)
    assert 'paragraphs' in result
    assert 'tables' in result
    assert len(result['tables']) == 1
    assert len(result['tables'][0]) == 2  # 2 rows
    assert len(result['tables'][0][0]) == 2  # 2 columns
    assert result['tables'][0][0][0] == "Cell 1"
    assert result['tables'][0][0][1] == "Cell 2"
    assert result['tables'][0][1][0] == "Cell 3"
    assert result['tables'][0][1][1] == "Cell 4"

@pytest.mark.asyncio
async def test_read_pptx(office_tool, sample_pptx_file):
    """Test reading content from a PPTX file."""
    result = await office_tool.run('read_pptx', file_path=sample_pptx_file)
    
    assert isinstance(result, list)
    assert len(result) >= 2
    assert "Test Slide" in result
    assert "This is a test slide." in result
    assert "Second Slide" in result
    assert "This is another test slide." in result

@pytest.mark.asyncio
async def test_read_xlsx(office_tool, sample_xlsx_file):
    """Test reading content from an XLSX file."""
    result = await office_tool.run('read_xlsx', file_path=sample_xlsx_file)
    
    assert isinstance(result, list)
    assert len(result) == 3
    assert result[0]['Name'] == 'Alice'
    assert result[0]['Age'] == 25
    assert result[0]['City'] == 'New York'
    assert result[1]['Name'] == 'Bob'
    assert result[2]['Name'] == 'Charlie'

@pytest.mark.asyncio
async def test_write_docx(office_tool, temp_dir):
    """Test writing content to a DOCX file."""
    output_path = os.path.join(temp_dir, "output.docx")
    result = await office_tool.run('write_docx', 
                                  text="This is a test document.\nIt has multiple lines.",
                                  output_path=output_path)
    
    assert isinstance(result, dict)
    assert 'success' in result
    assert result['success'] is True
    assert 'file_path' in result
    assert os.path.exists(output_path)
    
    # Verify the content by reading it back
    read_result = await office_tool.run('read_docx', file_path=output_path)
    assert "This is a test document." in read_result['paragraphs']

@pytest.mark.asyncio
async def test_write_docx_with_table(office_tool, temp_dir):
    """Test writing content to a DOCX file with a table."""
    output_path = os.path.join(temp_dir, "output_with_table.docx")
    result = await office_tool.run('write_docx', 
                                  text="This is a test document with a table.",
                                  output_path=output_path,
                                  table_data=[["Header 1", "Header 2"], ["Value 1", "Value 2"]])
    
    assert isinstance(result, dict)
    assert 'success' in result
    assert result['success'] is True
    assert 'file_path' in result
    assert os.path.exists(output_path)
    
    # Verify the content by reading it back
    read_result = await office_tool.run('read_docx', file_path=output_path, include_tables=True)
    assert "This is a test document with a table." in read_result['paragraphs']
    assert len(read_result['tables']) == 1
    assert read_result['tables'][0][0][0] == "Header 1"
    assert read_result['tables'][0][0][1] == "Header 2"
    assert read_result['tables'][0][1][0] == "Value 1"
    assert read_result['tables'][0][1][1] == "Value 2"

@pytest.mark.asyncio
async def test_write_pptx(office_tool, temp_dir):
    """Test writing content to a PPTX file."""
    output_path = os.path.join(temp_dir, "output.pptx")
    result = await office_tool.run('write_pptx', 
                                  slides=["Slide 1 Content", "Slide 2 Content"],
                                  output_path=output_path)
    
    assert isinstance(result, dict)
    assert 'success' in result
    assert result['success'] is True
    assert 'file_path' in result
    assert os.path.exists(output_path)
    
    # Verify the content by reading it back
    read_result = await office_tool.run('read_pptx', file_path=output_path)
    assert "Slide 1 Content" in " ".join(read_result)
    assert "Slide 2 Content" in " ".join(read_result)

@pytest.mark.asyncio
async def test_write_xlsx(office_tool, temp_dir):
    """Test writing content to an XLSX file."""
    output_path = os.path.join(temp_dir, "output.xlsx")
    data = [
        {"Name": "Alice", "Age": 25, "City": "New York"},
        {"Name": "Bob", "Age": 30, "City": "London"},
        {"Name": "Charlie", "Age": 35, "City": "Paris"}
    ]
    result = await office_tool.run('write_xlsx', 
                                  data=data,
                                  output_path=output_path,
                                  sheet_name="People")
    
    assert isinstance(result, dict)
    assert 'success' in result
    assert result['success'] is True
    assert 'file_path' in result
    assert os.path.exists(output_path)
    
    # Verify the content by reading it back
    read_result = await office_tool.run('read_xlsx', file_path=output_path)
    assert len(read_result) == 3
    assert read_result[0]['Name'] == 'Alice'
    assert read_result[1]['Name'] == 'Bob'
    assert read_result[2]['Name'] == 'Charlie'

@pytest.mark.asyncio
async def test_extract_text(office_tool, sample_pdf_file):
    """Test extracting text from a PDF file."""
    # Mock textract.process to return the expected text
    with patch('app.tools.office_tool.textract.process') as mock_process:
        mock_process.return_value = b"This is extracted text from the PDF."
        
        result = await office_tool.run('extract_text', file_path=sample_pdf_file)
        
        assert isinstance(result, str)
        assert result == "This is extracted text from the PDF."
        mock_process.assert_called_once_with(sample_pdf_file, encoding='utf-8')

# Error handling tests
@pytest.mark.asyncio
async def test_invalid_operation(office_tool):
    """Test handling of invalid operations."""
    with pytest.raises(OfficeToolError):
        await office_tool.run('invalid_op', file_path='some_file.docx')

@pytest.mark.asyncio
async def test_file_not_found(office_tool):
    """Test handling of non-existent files."""
    with pytest.raises(FileOperationError):
        await office_tool.run('read_docx', file_path='nonexistent_file.docx')

@pytest.mark.asyncio
async def test_invalid_extension(office_tool, temp_dir):
    """Test handling of invalid file extensions."""
    invalid_file = os.path.join(temp_dir, "invalid.txt")
    with open(invalid_file, 'w') as f:
        f.write("This is not a valid office file")
    
    with pytest.raises(SecurityError):
        await office_tool.run('read_docx', file_path=invalid_file)

@pytest.mark.asyncio
async def test_path_traversal_attempt(office_tool):
    """Test handling of path traversal attempts."""
    with pytest.raises(SecurityError):
        await office_tool.run('read_docx', file_path='../../../etc/passwd')

@pytest.mark.asyncio
async def test_content_validation_error(office_tool, temp_dir):
    """Test handling of invalid document content."""
    invalid_docx = os.path.join(temp_dir, "invalid.docx")
    with open(invalid_docx, 'wb') as f:
        f.write(b"This is not a valid DOCX file")
    
    # Mock validate_document to raise ContentValidationError
    with patch('app.tools.office_tool.validate_document') as mock_validate:
        mock_validate.side_effect = ContentValidationError("Invalid DOCX structure")
        
        with pytest.raises(ContentValidationError):
            await office_tool.run('read_docx', file_path=invalid_docx)

@pytest.mark.asyncio
async def test_cache_functionality(office_tool, sample_docx_file):
    """Test that caching works correctly."""
    # First call should not be cached
    result1 = await office_tool.run('read_docx', file_path=sample_docx_file)
    
    # Mock the validate_document function to verify it's not called on second run
    with patch('app.tools.office_tool.validate_document') as mock_validate:
        # Second call should use cache
        result2 = await office_tool.run('read_docx', file_path=sample_docx_file)
        
        # Validate document should not be called for cached result
        mock_validate.assert_not_called()
    
    assert result1 == result2, "Cached result should be identical to original"

@pytest.mark.asyncio
async def test_sanitize_text():
    """Test text sanitization in write_docx."""
    # Create a string with potentially harmful control characters
    unsafe_text = "Normal text\x00with\x1Fcontrol\x7Fcharacters"
    
    # Import the sanitize_text function from the module
    from app.tools.office_tool import write_docx
    
    # Access the nested sanitize_text function
    # This is a bit hacky but allows testing the nested function
    sanitize_text = None
    for name, func in write_docx.__globals__.items():
        if name == 'sanitize_text':
            sanitize_text = func
            break
    
    if sanitize_text:
        result = sanitize_text(unsafe_text)
        assert "\x00" not in result
        assert "\x1F" not in result
        assert "\x7F" not in result
        assert "Normal text" in result
        assert "withcontrolcharacters" in result
    else:
        pytest.skip("Could not access sanitize_text function")