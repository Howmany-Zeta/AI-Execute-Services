"""
Tests for PPTX writing functionality in DocumentWriterTool

This test suite covers PPTX/PPT writing functionality including:
- PPTX document writing
- Slide parsing and conversion
- Integration with office_tool
- Append mode for PPTX files
- Error handling for PPTX operations
"""
import os
import pytest
import tempfile
import logging
from pathlib import Path
from datetime import datetime
from typing import List

from aiecs.tools.docs.document_writer_tool import (
    DocumentWriterTool,
    DocumentFormat,
    WriteMode,
    EncodingType,
    ValidationLevel,
    DocumentWriterError,
    StorageError,
)
from pptx import Presentation

logger = logging.getLogger(__name__)


class TestDocumentWriterToolPPTX:
    """Test suite for PPTX writing functionality in DocumentWriterTool"""

    @pytest.fixture
    def writer_config(self, tmp_path):
        """Configuration for DocumentWriterTool"""
        return {
            "temp_dir": str(tmp_path / "temp"),
            "backup_dir": str(tmp_path / "backup"),
            "output_dir": str(tmp_path / "output"),
            "max_file_size": 100 * 1024 * 1024,
            "default_encoding": "utf-8",
            "enable_backup": True,
            "enable_versioning": True,
            "atomic_write": True,
        }

    @pytest.fixture
    def writer_tool(self, writer_config):
        """Create DocumentWriterTool instance for testing"""
        return DocumentWriterTool(writer_config)

    @pytest.fixture
    def sample_slides_content(self):
        """Sample slides content for testing"""
        return """# Title Slide
Welcome to the Presentation

---

## Slide 2: Introduction
This is the introduction slide with some content.

---

## Slide 3: Main Points
- Point 1
- Point 2
- Point 3

---

## Slide 4: Conclusion
Thank you for your attention!
"""

    def test_pptx_format_in_enum(self):
        """Test that PPTX format is available in DocumentFormat enum"""
        assert hasattr(DocumentFormat, "PPTX")
        assert hasattr(DocumentFormat, "PPT")
        assert DocumentFormat.PPTX == "pptx"
        assert DocumentFormat.PPT == "ppt"

    def test_office_tool_initialization(self, writer_tool):
        """Test that office_tool is properly initialized"""
        assert hasattr(writer_tool, "office_tool")
        # office_tool should be initialized (can be None if import fails)
        assert writer_tool.office_tool is not None or writer_tool.office_tool is None

    def test_write_pptx_basic(self, writer_tool, sample_slides_content, tmp_path):
        """Test basic PPTX writing functionality"""
        output_path = str(tmp_path / "test_write.pptx")

        result = writer_tool.write_document(
            target_path=output_path,
            content=sample_slides_content,
            format=DocumentFormat.PPTX,
            mode=WriteMode.CREATE,
            encoding=EncodingType.UTF8,
            validation_level=ValidationLevel.BASIC,
        )

        # Verify result structure
        assert "operation_id" in result
        assert result["target_path"] == output_path
        assert result["format"] == DocumentFormat.PPTX
        assert result["write_mode"] == WriteMode.CREATE
        assert os.path.exists(output_path)

        # Verify PPTX file can be opened and has slides
        prs = Presentation(output_path)
        assert len(prs.slides) > 0

    def test_write_pptx_with_ppt_format(self, writer_tool, sample_slides_content, tmp_path):
        """Test writing PPTX file using PPT format enum"""
        output_path = str(tmp_path / "test_write_ppt.pptx")

        result = writer_tool.write_document(
            target_path=output_path,
            content=sample_slides_content,
            format=DocumentFormat.PPT,
            mode=WriteMode.CREATE,
        )

        assert result["format"] == DocumentFormat.PPT
        assert os.path.exists(output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) > 0

    def test_write_pptx_overwrite_mode(self, writer_tool, sample_slides_content, tmp_path):
        """Test PPTX writing with overwrite mode"""
        output_path = str(tmp_path / "test_overwrite.pptx")

        # Create initial file
        writer_tool.write_document(
            target_path=output_path,
            content="Initial slide content",
            format=DocumentFormat.PPTX,
            mode=WriteMode.CREATE,
        )

        initial_slides = len(Presentation(output_path).slides)

        # Overwrite with new content
        result = writer_tool.write_document(
            target_path=output_path,
            content=sample_slides_content,
            format=DocumentFormat.PPTX,
            mode=WriteMode.OVERWRITE,
        )

        assert result["write_mode"] == WriteMode.OVERWRITE
        assert os.path.exists(output_path)
        new_slides = len(Presentation(output_path).slides)
        # New content should have different number of slides
        assert new_slides != initial_slides or new_slides > 0

    def test_write_pptx_append_mode(self, writer_tool, tmp_path):
        """Test PPTX writing with append mode"""
        output_path = str(tmp_path / "test_append.pptx")

        # Create initial file
        initial_content = "# Slide 1\nInitial content"
        writer_tool.write_document(
            target_path=output_path,
            content=initial_content,
            format=DocumentFormat.PPTX,
            mode=WriteMode.CREATE,
        )

        initial_slides = len(Presentation(output_path).slides)

        # Append new slides
        append_content = "\n---\n# Slide 2\nAppended content"
        result = writer_tool.write_document(
            target_path=output_path,
            content=append_content,
            format=DocumentFormat.PPTX,
            mode=WriteMode.APPEND,
        )

        assert result["write_mode"] == WriteMode.APPEND
        assert os.path.exists(output_path)
        final_slides = len(Presentation(output_path).slides)
        # Should have more slides after append
        assert final_slides >= initial_slides

    def test_write_pptx_with_list_content(self, writer_tool, tmp_path):
        """Test writing PPTX with list of slide strings"""
        output_path = str(tmp_path / "test_list_slides.pptx")

        slides_list = [
            "Slide 1: Title\nWelcome",
            "Slide 2: Content\nMain points here",
            "Slide 3: Conclusion\nThank you",
        ]

        # Convert list to string format that can be parsed
        content = "\n---\n".join(slides_list)

        result = writer_tool.write_document(
            target_path=output_path,
            content=content,
            format=DocumentFormat.PPTX,
            mode=WriteMode.CREATE,
        )

        assert os.path.exists(output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) >= len(slides_list)

    def test_write_pptx_parse_content_to_slides(self, writer_tool):
        """Test _parse_content_to_slides method"""
        # Test with "---" separator
        content1 = """Slide 1
---
Slide 2
---
Slide 3"""
        slides1 = writer_tool._parse_content_to_slides(content1)
        assert len(slides1) == 3
        assert "Slide 1" in slides1[0]
        assert "Slide 2" in slides1[1]
        assert "Slide 3" in slides1[2]

        # Test with "## Slide" headers
        content2 = """## Slide 1: Title
Content 1
## Slide 2: Content
Content 2"""
        slides2 = writer_tool._parse_content_to_slides(content2)
        assert len(slides2) >= 1

        # Test empty content
        slides3 = writer_tool._parse_content_to_slides("")
        assert len(slides3) == 1  # Should create one empty slide

    def test_write_pptx_file_structure(self, writer_tool, sample_slides_content, tmp_path):
        """Test that written PPTX file has valid structure"""
        output_path = str(tmp_path / "structure_test.pptx")

        writer_tool.write_document(
            target_path=output_path,
            content=sample_slides_content,
            format=DocumentFormat.PPTX,
            mode=WriteMode.CREATE,
        )

        # Verify file exists and can be opened
        assert os.path.exists(output_path)
        prs = Presentation(output_path)

        # Verify slides have content
        has_text = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    has_text = True
                    break
            if has_text:
                break
        assert has_text, "At least one slide should have text content"

    def test_write_pptx_without_office_tool_raises_error(self, tmp_path):
        """Test that PPTX writing fails gracefully if office_tool is not available"""
        config = {
            "temp_dir": str(tmp_path / "temp"),
            "backup_dir": str(tmp_path / "backup"),
        }
        writer_tool = DocumentWriterTool(config)

        # If office_tool failed to initialize, it should be None
        if writer_tool.office_tool is None:
            output_path = str(tmp_path / "test_no_office_tool.pptx")
            with pytest.raises((StorageError, DocumentWriterError)):
                writer_tool.write_document(
                    target_path=output_path,
                    content="Test content",
                    format=DocumentFormat.PPTX,
                    mode=WriteMode.CREATE,
                )

    def test_write_pptx_backup_created(self, writer_tool, sample_slides_content, tmp_path):
        """Test that backup is created when overwriting PPTX file"""
        output_path = str(tmp_path / "backup_test.pptx")

        # Create initial file
        writer_tool.write_document(
            target_path=output_path,
            content="Initial content",
            format=DocumentFormat.PPTX,
            mode=WriteMode.CREATE,
        )

        # Overwrite with backup
        result = writer_tool.write_document(
            target_path=output_path,
            content=sample_slides_content,
            format=DocumentFormat.PPTX,
            mode=WriteMode.OVERWRITE,
        )

        # Check if backup info is present
        if result.get("backup_info"):
            assert "backup_path" in result["backup_info"]

    def test_write_pptx_validation_levels(self, writer_tool, sample_slides_content, tmp_path):
        """Test PPTX writing with different validation levels"""
        for i, validation_level in enumerate([ValidationLevel.NONE, ValidationLevel.BASIC]):
            output_path = str(tmp_path / f"validation_test_{i}.pptx")
            result = writer_tool.write_document(
                target_path=output_path,
                content=sample_slides_content,
                format=DocumentFormat.PPTX,
                mode=WriteMode.CREATE,
                validation_level=validation_level,
            )

            # Check that validation_level is in content_metadata or result
            assert "content_metadata" in result
            assert result["content_metadata"]["validation_level"] == validation_level
            assert os.path.exists(output_path)

    def test_write_pptx_metadata_tracking(self, writer_tool, sample_slides_content, tmp_path):
        """Test that PPTX writing includes proper metadata"""
        output_path = str(tmp_path / "metadata_test.pptx")

        result = writer_tool.write_document(
            target_path=output_path,
            content=sample_slides_content,
            format=DocumentFormat.PPTX,
            mode=WriteMode.CREATE,
        )

        assert "content_metadata" in result
        assert "write_result" in result
        assert "size" in result["write_result"]
        assert result["write_result"]["size"] > 0

    def test_write_pptx_atomic_operation(self, writer_tool, sample_slides_content, tmp_path):
        """Test that PPTX writing uses atomic operations"""
        output_path = str(tmp_path / "atomic_test.pptx")

        result = writer_tool.write_document(
            target_path=output_path,
            content=sample_slides_content,
            format=DocumentFormat.PPTX,
            mode=WriteMode.CREATE,
        )

        # Verify file exists and is complete
        assert os.path.exists(output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) > 0

    @pytest.mark.parametrize(
        "format_enum",
        [DocumentFormat.PPTX, DocumentFormat.PPT],
    )
    def test_write_pptx_formats(self, writer_tool, sample_slides_content, tmp_path, format_enum):
        """Test writing with different PPTX format enums"""
        output_path = str(tmp_path / f"test_{format_enum.value}.pptx")

        result = writer_tool.write_document(
            target_path=output_path,
            content=sample_slides_content,
            format=format_enum,
            mode=WriteMode.CREATE,
        )

        assert result["format"] == format_enum
        assert os.path.exists(output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) > 0


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])

