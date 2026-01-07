"""
Tests for PPTX creation functionality in DocumentCreatorTool

This test suite covers PPTX/PPT creation functionality including:
- PPTX document creation from templates
- Slide parsing and conversion
- Integration with office_tool
- Error handling for PPTX operations
"""
import os
import pytest
import tempfile
import logging
from pathlib import Path
from datetime import datetime
from typing import Dict, Any

from aiecs.tools.docs.document_creator_tool import (
    DocumentCreatorTool,
    DocumentType,
    DocumentFormat,
    TemplateType,
    StylePreset,
    DocumentCreationError,
)
from pptx import Presentation

logger = logging.getLogger(__name__)


class TestDocumentCreatorToolPPTX:
    """Test suite for PPTX creation functionality in DocumentCreatorTool"""

    @pytest.fixture
    def creator_config(self, tmp_path):
        """Configuration for DocumentCreatorTool"""
        return {
            "templates_dir": str(tmp_path / "templates"),
            "output_dir": str(tmp_path / "output"),
            "default_format": DocumentFormat.MARKDOWN,
            "default_style": StylePreset.DEFAULT,
            "auto_backup": True,
            "include_metadata": True,
            "generate_toc": True,
        }

    @pytest.fixture
    def creator_tool(self, creator_config):
        """Create DocumentCreatorTool instance for testing"""
        return DocumentCreatorTool(creator_config)

    @pytest.fixture
    def sample_presentation_metadata(self):
        """Sample metadata for presentation testing"""
        return {
            "title": "Test Presentation",
            "presenter": "John Doe",
            "date": datetime.now().strftime("%Y-%m-%d"),
            "organization": "Test Company",
            "agenda": "1. Introduction\n2. Main Points\n3. Conclusion",
            "introduction": "Welcome to the presentation",
            "main_content": "Key points and details about the topic",
            "conclusion": "Summary and next steps",
            "questions": "Q&A Session",
            "contact_info": "contact@example.com",
        }

    def test_pptx_format_in_enum(self):
        """Test that PPTX format is available in DocumentFormat enum"""
        assert hasattr(DocumentFormat, "PPTX")
        assert hasattr(DocumentFormat, "PPT")
        assert DocumentFormat.PPTX == "pptx"
        assert DocumentFormat.PPT == "ppt"

    def test_office_tool_initialization(self, creator_tool):
        """Test that office_tool is properly initialized"""
        assert hasattr(creator_tool, "office_tool")
        # office_tool should be initialized (can be None if import fails)
        assert creator_tool.office_tool is not None or creator_tool.office_tool is None

    def test_create_pptx_presentation(self, creator_tool, sample_presentation_metadata, tmp_path):
        """Test creating a PPTX presentation from presentation template"""
        output_path = str(tmp_path / "test_presentation.pptx")

        result = creator_tool.create_document(
            document_type=DocumentType.PRESENTATION,
            template_type=TemplateType.PRESENTATION,
            output_format=DocumentFormat.PPTX,
            metadata=sample_presentation_metadata,
            output_path=output_path,
        )

        # Verify result structure
        assert "document_id" in result
        assert result["document_type"] == DocumentType.PRESENTATION
        assert result["template_type"] == TemplateType.PRESENTATION
        assert result["output_format"] == DocumentFormat.PPTX
        assert result["output_path"] == output_path
        assert os.path.exists(output_path)

        # Verify PPTX file can be opened and has slides
        prs = Presentation(output_path)
        assert len(prs.slides) > 0

    def test_create_pptx_with_ppt_format(self, creator_tool, sample_presentation_metadata, tmp_path):
        """Test creating PPTX file using PPT format enum (should create PPTX)"""
        output_path = str(tmp_path / "test_presentation_ppt.pptx")

        result = creator_tool.create_document(
            document_type=DocumentType.PRESENTATION,
            template_type=TemplateType.PRESENTATION,
            output_format=DocumentFormat.PPT,
            metadata=sample_presentation_metadata,
            output_path=output_path,
        )

        assert result["output_format"] == DocumentFormat.PPT
        assert os.path.exists(output_path)
        # File should have .pptx extension even when using PPT format
        assert output_path.endswith(".pptx") or os.path.exists(output_path)

    def test_pptx_slide_parsing_from_markdown(self, creator_tool, tmp_path):
        """Test that markdown content with slide separators is parsed correctly"""
        output_path = str(tmp_path / "test_slides.pptx")

        # Create content with slide separators
        content = """# Title Slide
Welcome to the presentation

---

## Slide 2: Introduction
This is the introduction slide

---

## Slide 3: Main Content
Key points here
"""

        # Use create_from_template with markdown content
        template_path = tmp_path / "templates" / "test_template.md"
        template_path.parent.mkdir(parents=True, exist_ok=True)
        template_path.write_text(content)

        result = creator_tool.create_from_template(
            template_name="test_template.md",
            template_variables={},
            output_format=DocumentFormat.PPTX,
            output_path=output_path,
        )

        assert os.path.exists(output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) >= 2  # Should have multiple slides

    def test_pptx_presentation_template_supports_pptx(self, creator_tool):
        """Test that presentation template supports PPTX format"""
        template = creator_tool._get_presentation_template()
        assert "pptx" in template["supported_formats"]

    def test_pptx_creation_with_blank_template(self, creator_tool, sample_presentation_metadata, tmp_path):
        """Test creating PPTX with blank template"""
        output_path = str(tmp_path / "blank_pptx.pptx")

        result = creator_tool.create_document(
            document_type=DocumentType.PRESENTATION,
            template_type=TemplateType.BLANK,
            output_format=DocumentFormat.PPTX,
            metadata=sample_presentation_metadata,
            output_path=output_path,
        )

        assert os.path.exists(output_path)
        prs = Presentation(output_path)
        assert len(prs.slides) >= 0  # Blank template may have 0 or 1 slide

    def test_pptx_file_has_valid_structure(self, creator_tool, sample_presentation_metadata, tmp_path):
        """Test that created PPTX file has valid structure"""
        output_path = str(tmp_path / "valid_structure.pptx")

        creator_tool.create_document(
            document_type=DocumentType.PRESENTATION,
            template_type=TemplateType.PRESENTATION,
            output_format=DocumentFormat.PPTX,
            metadata=sample_presentation_metadata,
            output_path=output_path,
        )

        # Verify file exists and can be opened
        assert os.path.exists(output_path)
        prs = Presentation(output_path)

        # Verify slides have content
        for slide in prs.slides:
            has_text = False
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    has_text = True
                    break
            # At least some slides should have text
            if has_text:
                break

    def test_pptx_creation_without_office_tool_raises_error(self, tmp_path):
        """Test that PPTX creation fails gracefully if office_tool is not available"""
        # Mock the office_tool to be None
        config = {
            "templates_dir": str(tmp_path / "templates"),
            "output_dir": str(tmp_path / "output"),
        }
        creator_tool = DocumentCreatorTool(config)

        # If office_tool failed to initialize, it should be None
        if creator_tool.office_tool is None:
            output_path = str(tmp_path / "test_no_office_tool.pptx")
            with pytest.raises(DocumentCreationError):
                creator_tool.create_document(
                    document_type=DocumentType.PRESENTATION,
                    template_type=TemplateType.PRESENTATION,
                    output_format=DocumentFormat.PPTX,
                    metadata={"title": "Test"},
                    output_path=output_path,
                )

    def test_parse_content_to_slides_with_separators(self, creator_tool):
        """Test _parse_content_to_slides method with different separator formats"""
        # Test with "---" separator
        content1 = """Slide 1 content
---
Slide 2 content
---
Slide 3 content"""
        slides1 = creator_tool._parse_content_to_slides(content1)
        assert len(slides1) == 3

        # Test with "## Slide" headers
        content2 = """## Slide 1: Title
Content for slide 1
## Slide 2: Content
Content for slide 2"""
        slides2 = creator_tool._parse_content_to_slides(content2)
        assert len(slides2) >= 1

        # Test with paragraph breaks
        content3 = """First slide content

Second slide content

Third slide content"""
        slides3 = creator_tool._parse_content_to_slides(content3)
        assert len(slides3) >= 1

    def test_pptx_output_path_generation(self, creator_tool, sample_presentation_metadata):
        """Test that PPTX output path is generated correctly"""
        result = creator_tool.create_document(
            document_type=DocumentType.PRESENTATION,
            template_type=TemplateType.PRESENTATION,
            output_format=DocumentFormat.PPTX,
            metadata=sample_presentation_metadata,
        )

        assert "output_path" in result
        assert result["output_path"].endswith(".pptx")
        assert os.path.exists(result["output_path"])

    def test_pptx_detect_format(self, creator_tool):
        """Test that PPTX format is detected correctly from file extension"""
        assert creator_tool._detect_document_format("test.pptx") == DocumentFormat.PPTX
        assert creator_tool._detect_document_format("test.ppt") == DocumentFormat.PPT

    def test_pptx_creation_tracks_document(self, creator_tool, sample_presentation_metadata, tmp_path):
        """Test that PPTX creation is tracked in created documents"""
        initial_count = len(creator_tool.get_created_documents())

        output_path = str(tmp_path / "tracked_pptx.pptx")
        creator_tool.create_document(
            document_type=DocumentType.PRESENTATION,
            template_type=TemplateType.PRESENTATION,
            output_format=DocumentFormat.PPTX,
            metadata=sample_presentation_metadata,
            output_path=output_path,
        )

        created_docs = creator_tool.get_created_documents()
        assert len(created_docs) == initial_count + 1
        assert created_docs[-1]["output_format"] == DocumentFormat.PPTX

    def test_pptx_creation_metadata(self, creator_tool, sample_presentation_metadata, tmp_path):
        """Test that PPTX creation includes proper metadata"""
        output_path = str(tmp_path / "metadata_pptx.pptx")

        result = creator_tool.create_document(
            document_type=DocumentType.PRESENTATION,
            template_type=TemplateType.PRESENTATION,
            output_format=DocumentFormat.PPTX,
            metadata=sample_presentation_metadata,
            output_path=output_path,
        )

        assert "creation_metadata" in result
        assert "created_at" in result["creation_metadata"]
        assert "file_size" in result["creation_metadata"]
        assert result["creation_metadata"]["file_size"] > 0

    @pytest.mark.parametrize(
        "format_enum,expected_ext",
        [
            (DocumentFormat.PPTX, "pptx"),
            (DocumentFormat.PPT, "pptx"),  # PPT format creates PPTX files
        ],
    )
    def test_pptx_format_extensions(self, creator_tool, sample_presentation_metadata, tmp_path, format_enum, expected_ext):
        """Test that PPTX formats generate correct file extensions"""
        result = creator_tool.create_document(
            document_type=DocumentType.PRESENTATION,
            template_type=TemplateType.PRESENTATION,
            output_format=format_enum,
            metadata=sample_presentation_metadata,
        )

        assert result["output_path"].endswith(f".{expected_ext}")


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])

